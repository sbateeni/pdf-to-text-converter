import streamlit as st
from pptx import Presentation
from fpdf import FPDF
import tempfile

def convert_pptx_to_pdf(pptx_file):
    """Convert PowerPoint file to PDF."""
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    presentation = Presentation(pptx_file)
    for slide in presentation.slides:
        pdf.add_page()
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.set_font("Arial", size=12)
                pdf.multi_cell(0, 10, shape.text)

    pdf_file_path = "output.pdf"
    pdf.output(pdf_file_path)
    return pdf_file_path

def show():
    st.title("📄 Convert PowerPoint to PDF")

    uploaded_file = st.file_uploader("📂 Choose a PowerPoint file", type=["pptx"])

    if uploaded_file:
        if st.button("📝 Convert PowerPoint to PDF"):
            try:
                pdf_file_path = convert_pptx_to_pdf(uploaded_file)
                st.success("✅ Conversion completed!")

                with open(pdf_file_path, "rb") as pdf_file:
                    st.download_button(
                        label="📥 Download PDF Document",
                        data=pdf_file.read(),
                        file_name="output.pdf",
                        mime="application/pdf"
                    )
            except Exception as e:
                st.error(f"❌ An error occurred during conversion: {str(e)}")

if __name__ == "__main__":
    show()